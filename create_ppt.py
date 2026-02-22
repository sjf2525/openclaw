#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import sys

def create_openclaw_presentation(output_path):
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "OpenClaw：你的终极AI助手框架"
    subtitle.text = "部署、配置、技能使用与Feishu集成完全指南"
    
    # Slide 2: Why OpenClaw
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "为什么选择OpenClaw？"
    text_frame = content.text_frame
    text_frame.clear()
    
    points = [
        "开源自由：完全控制，无供应商锁定",
        "本地部署：数据不出域，隐私安全", 
        "技能生态：丰富插件，无限扩展",
        "多通道集成：Feishu、WhatsApp、Telegram一站式打通"
    ]
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(24)
    
    # Slide 3: Deployment steps
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "部署：三步启动"
    text_frame = content.text_frame
    text_frame.clear()
    
    steps = [
        "1. 运行Ollama - 本地AI引擎",
        "2. 配置Kimi K2.5 - 强大模型支持", 
        "3. 启动OpenClaw - 浏览器即刻访问"
    ]
    
    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(28)
    
    # Slide 4: Skills
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "技能：即插即用"
    text_frame = content.text_frame
    text_frame.clear()
    
    skills = [
        "office-document-specialist-suite：专业Office文档处理",
        "pdf-extract：PDF文本一键提取",
        "ppt-generator：乔布斯风演示稿自动生成",
        "浏览器控制：网页自动化全掌控"
    ]
    
    for skill in skills:
        p = text_frame.add_paragraph()
        p.text = skill
        p.level = 0
        p.font.size = Pt(24)
    
    # Slide 5: Feishu integration
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Feishu集成：团队协作革命"
    text_frame = content.text_frame
    text_frame.clear()
    
    feishu_points = [
        "创建应用：Feishu开放平台三步走",
        "配置命令：简单命令行设置",
        "桥接技能：feishu-bridge工作原理"
    ]
    
    for point in feishu_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(26)
    
    # Slide 6: Configuration commands
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "配置命令展示"
    text_frame = content.text_frame
    text_frame.clear()
    
    # Add code block as text
    code = """openclaw config set channels.feishu.appId "你的ID"
openclaw config set channels.feishu.appSecret "你的密钥"
openclaw config set channels.feishu.enabled true"""
    
    p = text_frame.add_paragraph()
    p.text = code
    p.font.name = "Courier New"
    p.font.size = Pt(18)
    
    # Slide 7: Skill development
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "技能开发：创造你的专属能力"
    text_frame = content.text_frame
    text_frame.clear()
    
    dev_points = [
        "技能结构：目录布局",
        "开发流程：四步法",
        "工具注册：@tool装饰器"
    ]
    
    for point in dev_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(26)
    
    # Slide 8: Best practices
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "最佳实践：高效使用指南"
    text_frame = content.text_frame
    text_frame.clear()
    
    practices = [
        "安装策略：精选技能，定期更新",
        "应用场景：文档自动化、消息集成、研究助手",
        "故障排除：状态检查、服务重启、日志查看"
    ]
    
    for practice in practices:
        p = text_frame.add_paragraph()
        p.text = practice
        p.level = 0
        p.font.size = Pt(26)
    
    # Slide 9: Resources
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "资源网络：持续成长"
    text_frame = content.text_frame
    text_frame.clear()
    
    resources = [
        "官方文档：docs.openclaw.ai",
        "技能市场：clawhub.com",
        "GitHub：github.com/openclaw/openclaw",
        "社区：Discord邀请链接"
    ]
    
    for resource in resources:
        p = text_frame.add_paragraph()
        p.text = resource
        p.level = 0
        p.font.size = Pt(26)
    
    # Slide 10: Future vision
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "未来展望"
    text_frame = content.text_frame
    text_frame.clear()
    
    vision = [
        "OpenClaw不仅是一个工具，更是AI民主化的桥梁",
        "让每个团队都能拥有定制化的AI助手",
        "释放人类创造力，专注真正重要的事情"
    ]
    
    for line in vision:
        p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(28)
    
    # Slide 11: Call to action
    slide_layout = prs.slide_layouts[0]  # Title slide again
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "开始你的OpenClaw之旅"
    subtitle.text = "重塑工作方式！"
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    output_file = "openclaw_guide.pptx"
    if len(sys.argv) > 1:
        output_file = sys.argv[1]
    
    create_openclaw_presentation(output_file)